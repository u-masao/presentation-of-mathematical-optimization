import logging
import tempfile
from pathlib import Path

import click
import markdown
import mlflow
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Mm, Pt

# office xml open の drawingML namespace
NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def log_artifact_from_message(message, filename):
    """
    テキストメッセージをファイルに書き込んで mlflow に登録する
    """
    with tempfile.TemporaryDirectory() as temp_dir:
        file_path = Path(temp_dir) / filename
        open(file_path, "w").write(message)
        mlflow.log_artifact(file_path)


def get_layout_by_name(prs, query):
    """
    レイアウト名からレイアウトを取得する
    ない場合は ValueError を投げる
    """

    # init logger
    logger = logging.getLogger(__name__)

    # レイアウト名一覧を取得
    names = [x.name for x in prs.slide_layouts]
    logger.info(f"layout names: {names}")

    # レイアウト名を検索(なければ ValueError)
    index = names.index(query)

    # レイアウトを取得
    layout = prs.slide_layouts[index]
    logger.info(f"placeholders: {len(layout.placeholders)}")
    logger.info(f"placeholders name: {[x.name for x in layout.placeholders]}")
    return layout


def set_position_mm(placeholder, left=None, top=None, height=None, width=None):
    """
    Placeholder の位置と大きさを修正する
    """
    if left:
        placeholder.left = Mm(left)
    if top:
        placeholder.top = Mm(top)
    if height:
        placeholder.height = Mm(height)
    if width:
        placeholder.width = Mm(width)


def configure_presentation(
    prs,
    slide_width_mm=338.67,
    slide_height_mm=190.5,
):
    """
    プレゼンテーションの設定
    スライドレイアウトの調整
    """

    # プレゼンテーションのサイズを変更
    prs.slide_width = Mm(slide_width_mm)
    prs.slide_height = Mm(slide_height_mm)

    # レイアウト調整(コンテンツ)
    content_layout = get_layout_by_name(prs, "Title and Content")
    set_position_mm(
        content_layout.placeholders[0], 20, 20, 30, slide_width_mm - 2 * 20
    )
    set_position_mm(
        content_layout.placeholders[1],
        20,
        60,
        slide_height_mm - 80,
        slide_width_mm - 2 * 20,
    )

    # レイアウト調整(タイトルスライド)
    title_layout = get_layout_by_name(prs, "Title Slide")
    set_position_mm(
        title_layout.placeholders[0], 20, 20, 30, slide_width_mm - 2 * 20
    )
    set_position_mm(
        title_layout.placeholders[1],
        20,
        60,
        slide_height_mm - 80,
        slide_width_mm - 2 * 20,
    )


def make_presentation(html):
    # init logger
    logger = logging.getLogger(__name__)
    prs = Presentation()
    configure_presentation(prs)

    # コンテキストを変数に保持
    context = {
        "h1": "",
        "h2": "",
        "h3": "",
    }

    # スライド毎にループ
    for slide_html in html.split("<hr />"):

        # デバッグ表示
        logger.info(f"slide html: {slide_html}")

        # 空のスライドをスキップ
        if len(slide_html) == 0:
            continue

        # html をパース
        soup = BeautifulSoup(slide_html, "html.parser")

        # スライドレベルを特定
        slide_level = None
        for tag in ["h1", "h2", "h3"]:
            tag_html = soup.find(tag)
            if tag_html:
                slide_level = tag
                context[tag] = tag_html.get_text()

        # デバッグ表示
        logger.info(f"{context=}")

        # パースしたテキストを保持する変数を作成
        slide_texts = {"title": "", "context": "", "body": ""}

        # レベルに応じてタイトルとボディを設定
        if slide_level:
            title_soup = soup.find(slide_level)
            slide_texts["title"] = title_soup.get_text()
            slide_texts["body"] = title_soup.find_next_siblings()
        else:
            slide_texts["body"] = [soup]

        # パンくずを作成
        if slide_level == "h1":
            slide_texts["context"] = ""
        elif slide_level == "h2":
            slide_texts["context"] = context["h1"]
        elif slide_level == "h3":
            slide_texts["context"] = context["h1"] + " > " + context["h2"]

        # スライドを追加
        if slide_texts["title"]:
            add_slide(prs, slide_texts, slide_html)
        else:
            logger.warning(f"skip add slide: {slide_texts}")

    return prs


def add_slide(prs, slide_texts, slide_html):
    # init logger
    logger = logging.getLogger(__name__)
    logger.info(f"{slide_texts=}")

    # レイアウトを取得
    layout = get_layout_by_name(prs, "Title and Content")

    # スライドを追加
    slide = prs.slides.add_slide(layout)

    # title を設定
    title = slide.shapes.title
    title.text = slide_texts["title"]
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # コンテンツを設定
    content = slide.shapes.placeholders[1]  # content
    draw_soup_to_placeholder(content, slide_texts["body"])
    for x in content.text_frame.paragraphs:
        x.font.size = Pt(30)

    # note にデバッグ情報を追加
    slide.notes_slide.notes_text_frame.text = str(slide_html)


def parse_li(tag, level=1, ul_ol="ol"):
    result = []
    for li_tag in tag.find_all("li", recursive=False):
        result.append((level, li_tag.get_text(strip=True)))
        nested = li_tag.find(ul_ol)
        if nested:
            result.extend(parse_li(nested, level + 1))
    return result


def remove_unnumbered_list(pg):
    # 既存の <a:buChar> 要素を検索して削除
    bu_char_elem = pg._element.pPr.find(qn("a:buChar"))
    if bu_char_elem is not None:
        pg._element.pPr.remove(bu_char_elem)


def replace_to_numbered_list(pg, style: str = "arabicPlain"):

    remove_unnumbered_list(pg)

    # 新しい <a:buAutoNum> 要素を作成して追加
    bu_auto_num_elem = pg._element.pPr.makeelement(
        etree.QName(NSMAP["a"], "buAutoNum"), nsmap=NSMAP
    )
    bu_auto_num_elem.set("type", style)
    pg._element.pPr.append(bu_auto_num_elem)


def replace_bu_to_regular(pg):

    remove_unnumbered_list(pg)

    # 新しい <a:buAutoNum> 要素を作成して追加
    bu_auto_num_elem = pg._element.pPr.makeelement(
        etree.QName(NSMAP["a"], "buNone"), nsmap=NSMAP
    )
    pg._element.pPr.append(bu_auto_num_elem)


def draw_soup_to_placeholder(ph, soups):
    # init logger
    logger = logging.getLogger(__name__)

    for soup in soups:
        logger.info(f"{soup=}")

        if soup.name == "p" and len(list(soup.children)) == 1:
            # 'p' の場合
            pg = ph.text_frame.add_paragraph()
            pg.text = soup.get_text()
            pg.level = 0
            replace_bu_to_regular(pg)

        elif soup.name == "ol":
            # 'ol' 番号ありリストの場合
            li_items = parse_li(soup, ul_ol="ol", level=0)
            logger.info(f"{li_items=}")
            for level, text in li_items:
                pg = ph.text_frame.add_paragraph()
                pg.text = text
                pg.level = level
                replace_to_numbered_list(pg)

        elif soup.name == "ul":
            # 'ul' 番号なしリストの場合
            li_items = parse_li(soup, ul_ol="ul", level=0)
            logger.info(f"{li_items=}")
            for level, text in li_items:
                pg = ph.text_frame.add_paragraph()
                pg.text = text
                pg.level = level


def convert_markdown_to_pptx(md_text):
    html = markdown.markdown(md_text)
    return make_presentation(html)


@click.command()
@click.argument("input_filepath", type=click.Path(exists=True))
@click.argument("output_filepath", type=click.Path())
def main(**kwargs):

    # init logger
    logger = logging.getLogger(__name__)
    logger.info(f"args: {kwargs}")
    mlflow.set_experiment("make_pptx")
    mlflow.start_run()
    mlflow.log_params({f"args.{k}": v for k, v in kwargs.items()})

    # load prompt
    md_text = open(kwargs["input_filepath"], "r").read()

    # convert
    presentation = convert_markdown_to_pptx(md_text)

    # save file
    presentation.save(kwargs["output_filepath"])

    # logging
    mlflow.log_artifact(kwargs["output_filepath"])
    mlflow.end_run()


if __name__ == "__main__":
    log_fmt = "%(asctime)s - %(name)s - %(levelname)s - %(message)s"
    logging.basicConfig(level=logging.INFO, format=log_fmt)
    load_dotenv()
    main()
